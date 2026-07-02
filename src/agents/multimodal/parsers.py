"""Structured document extraction helpers for the FurnaceMind MRAG pipeline.

The public ``parse_*`` functions are kept for compatibility with older callers.
New MRAG ingestion uses the ``extract_*`` helpers so page, slide, sheet, and
image metadata can be carried into Qdrant payloads.
"""

from __future__ import annotations

import shutil
import subprocess
import tempfile
from io import BytesIO
from pathlib import Path
from typing import Any

import pandas as pd


def _read_bytes(file_or_bytes: Any) -> bytes:
    """Return bytes from a Streamlit upload, file-like object, or bytes value."""
    if isinstance(file_or_bytes, bytes):
        return file_or_bytes
    if isinstance(file_or_bytes, bytearray):
        return bytes(file_or_bytes)
    if hasattr(file_or_bytes, "seek"):
        try:
            file_or_bytes.seek(0)
        except Exception:
            pass
    data = file_or_bytes.read()
    if hasattr(file_or_bytes, "seek"):
        try:
            file_or_bytes.seek(0)
        except Exception:
            pass
    return data


def _df_to_markdown(df: pd.DataFrame) -> str:
    """Render a DataFrame as Markdown, falling back to CSV when tabulate is absent."""
    try:
        return df.to_markdown(index=False)
    except Exception:
        return df.to_csv(index=False)


def extract_pdf_pages(file_or_bytes: Any, *, render_pages: bool = True) -> list[dict]:
    """Extract text and optional rendered page images from a PDF."""
    import fitz  # PyMuPDF

    file_bytes = _read_bytes(file_or_bytes)
    doc = fitz.open(stream=file_bytes, filetype="pdf")
    pages: list[dict] = []
    for index, page in enumerate(doc, start=1):
        image_bytes: bytes | None = None
        if render_pages:
            pix = page.get_pixmap(matrix=fitz.Matrix(1.5, 1.5), alpha=False)
            image_bytes = pix.tobytes("png")
        pages.append(
            {
                "page_number": index,
                "text": page.get_text("text").strip(),
                "image_bytes": image_bytes,
            }
        )
    return pages


def extract_docx_text(file_or_bytes: Any) -> str:
    """Extract paragraph and table text from a DOCX file."""
    import docx

    document = docx.Document(BytesIO(_read_bytes(file_or_bytes)))
    parts: list[str] = []
    parts.extend(paragraph.text for paragraph in document.paragraphs if paragraph.text)
    for table in document.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            if any(cells):
                parts.append(" | ".join(cells))
    return "\n".join(part for part in parts if part).strip()


def extract_pptx_slides(file_or_bytes: Any) -> list[dict]:
    """Extract slide text and embedded image blobs from a PPTX file."""
    import pptx

    presentation = pptx.Presentation(BytesIO(_read_bytes(file_or_bytes)))
    slides: list[dict] = []
    for slide_index, slide in enumerate(presentation.slides, start=1):
        text_parts: list[str] = []
        image_blobs: list[bytes] = []
        for shape in slide.shapes:
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    if any(cells):
                        text_parts.append(" | ".join(cells))
            shape_text = getattr(shape, "text", "")
            if shape_text:
                text_parts.append(shape_text.strip())
            image = getattr(shape, "image", None)
            if image is not None:
                try:
                    image_blobs.append(image.blob)
                except Exception:
                    continue
        slides.append(
            {
                "slide_number": slide_index,
                "text": "\n".join(part for part in text_parts if part).strip(),
                "image_blobs": image_blobs,
            }
        )
    return slides


def render_pptx_slides(file_or_bytes: Any) -> list[dict]:
    """Render PPTX slides to PNG images when LibreOffice is available.

    python-pptx can extract text and embedded pictures but cannot render a full
    slide. LibreOffice gives us a best-effort headless conversion path so charts,
    layouts, and text-as-shape visuals can also be embedded by the multimodal
    model. If LibreOffice is not installed, callers simply continue with the
    extracted text and embedded image blobs.
    """
    executable = shutil.which("soffice") or shutil.which("libreoffice")
    if not executable:
        return []

    file_bytes = _read_bytes(file_or_bytes)
    with tempfile.TemporaryDirectory(prefix="fm_pptx_render_") as tmp:
        tmp_path = Path(tmp)
        input_path = tmp_path / "upload.pptx"
        output_dir = tmp_path / "out"
        output_dir.mkdir(parents=True, exist_ok=True)
        input_path.write_bytes(file_bytes)

        command = [
            executable,
            "--headless",
            "--convert-to",
            "pdf",
            "--outdir",
            str(output_dir),
            str(input_path),
        ]
        try:
            completed = subprocess.run(
                command,
                capture_output=True,
                check=False,
                timeout=90,
            )
        except Exception:
            return []
        if completed.returncode != 0:
            return []

        pdf_files = sorted(output_dir.glob("*.pdf"))
        if not pdf_files:
            return []

        import fitz  # PyMuPDF

        rendered: list[dict] = []
        doc = fitz.open(pdf_files[0])
        for index, page in enumerate(doc, start=1):
            pix = page.get_pixmap(matrix=fitz.Matrix(1.5, 1.5), alpha=False)
            rendered.append(
                {
                    "slide_number": index,
                    "image_bytes": pix.tobytes("png"),
                }
            )
        return rendered


def extract_excel_sheets(file_or_bytes: Any) -> list[dict]:
    """Extract every Excel sheet as Markdown table text."""
    workbook = pd.ExcelFile(BytesIO(_read_bytes(file_or_bytes)))
    sheets: list[dict] = []
    for sheet_name in workbook.sheet_names:
        df = pd.read_excel(workbook, sheet_name=sheet_name)
        sheets.append(
            {
                "sheet_name": sheet_name,
                "text": _df_to_markdown(df).strip(),
            }
        )
    return sheets


def extract_csv_text(file_or_bytes: Any) -> str:
    """Extract a CSV file as Markdown table text."""
    df = pd.read_csv(BytesIO(_read_bytes(file_or_bytes)))
    return _df_to_markdown(df).strip()


def extract_text_file(file_or_bytes: Any) -> str:
    """Decode a plain text or Markdown file."""
    return _read_bytes(file_or_bytes).decode("utf-8", errors="replace").strip()


def parse_pdf(file) -> str:
    """Extract all text from a PDF file."""
    return "\n".join(
        page["text"] for page in extract_pdf_pages(file, render_pages=False)
    ).strip()


def parse_docx(file) -> str:
    """Extract paragraph and table text from a DOCX file."""
    return extract_docx_text(file)


def parse_pptx(file) -> str:
    """Extract text from all slides in a PPTX file."""
    return "\n".join(slide["text"] for slide in extract_pptx_slides(file)).strip()


def parse_excel(file) -> str:
    """Read all Excel sheets and render them as Markdown tables."""
    return "\n\n".join(
        f"Sheet: {sheet['sheet_name']}\n{sheet['text']}"
        for sheet in extract_excel_sheets(file)
    ).strip()
