import pandas as pd
import docx
import pptx
import fitz  # PyMuPDF
from io import BytesIO


def parse_pdf(file):
    doc = fitz.open(stream=file.read(), filetype="pdf")
    return "\n".join([page.get_text() for page in doc])


def parse_docx(file):
    doc = docx.Document(file)
    return "\n".join([para.text for para in doc.paragraphs])


def parse_pptx(file):
    prs = pptx.Presentation(file)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text)


def parse_excel(file):
    df = pd.read_excel(file)
    return df.to_markdown()